import os
from fastapi import FastAPI, File, UploadFile, Form, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
from typing_extensions import Annotated, TypedDict, List
from typing import List, Optional
from dotenv import load_dotenv
import PyPDF2
import pptx  # For PPTX extraction
import docx  # For DOCX extraction
import csv   # For CSV extraction
import io
from io import BytesIO
from langchain_groq import ChatGroq
from langchain.prompts import ChatPromptTemplate
from langchain.text_splitter import RecursiveCharacterTextSplitter

# Load environment variables
load_dotenv()
api_key = os.getenv("GROQ_API_KEY")

# Set up FastAPI
app = FastAPI()

# Enable CORS for all origins, can be restricted in production
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

class Flashcard(TypedDict):
    """Flashcard for learning."""
    question: Annotated[str, "The question or prompt on the front of the flashcard"]
    answer: Annotated[str, "The answer or explanation on the back of the flashcard"]

# A set of flashcards
class FlashcardSet(TypedDict):
    """Set of flashcards."""
    flashcards: Annotated[List[Flashcard], "A list of flashcards based on the input text"]

# Define the Cloze Deletion Flashcard structure
class ClozeDeletionFlashcard(TypedDict):
    """Flashcard for cloze deletion."""
    question_with_blanks: Annotated[str, "A sentence with one or more blanks (____) for the learner to fill in."]
    correct_answers: Annotated[List[str], "The list of correct words or phrases that fill in the blanks."]

# A set of cloze deletion flashcards
class ClozeDeletionFlashcardSet(TypedDict):
    """Set of cloze deletion flashcards."""
    flashcards: Annotated[List[ClozeDeletionFlashcard], "A list of fill-in-the-blank flashcards based on the input text."]

# Text extraction functions
def extract_text_from_pdf(file) -> str:
    reader = PyPDF2.PdfReader(file)
    text = ""
    for page in reader.pages:
        text += page.extract_text() or ""
    return text

def extract_text_from_pptx(file) -> str:
    file_bytes = BytesIO(file.read())
    presentation = pptx.Presentation(file_bytes)
    text = ""
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    file.seek(0)
    return text

def extract_text_from_docx(file) -> str:
    file_bytes = BytesIO(file.read())
    doc = docx.Document(file_bytes)
    text = ""
    for paragraph in doc.paragraphs:
        text += paragraph.text + "\n"
    file.seek(0)
    return text

def extract_text_from_csv(file) -> str:
    decoded_file = io.StringIO(file.read().decode("utf-8"))
    text = ""
    reader = csv.reader(decoded_file)
    for row in reader:
        text += " ".join(row) + "\n"
    file.seek(0)
    return text

# Set up the LLM
llm = ChatGroq(temperature=0, model="llama3-8b-8192")

# Define prompts for different flashcard types
normal_prompt = ChatPromptTemplate.from_messages([
    ("system", "Generate a set of flashcards from the given text. Focus on key concepts and important information."),
    ("human", "Text: {chunk}\n\nGenerate a list of flashcards based on this text.")
])

cloze_prompt = ChatPromptTemplate.from_messages([
    ("system", 
     "Generate a set of cloze deletion flashcards from the given text. "
     "Focus on creating simple fill-in-the-blank questions for key facts or important information."),
    ("human", 
     "Text: {chunk}\n\nCreate flashcards using simple cloze deletions. For each flashcard, replace key information with blanks (____) "
     "to test the learner's memory. Provide the correct word or phrase for each blank in the answer.")
])

@app.get("/")
def read_root():
    return {"message": "Welcome to the flashcard generation prototype!"}

@app.post("/flashcard/")
async def create_flashcards(
    type: str = Form(...),
    method: str = Form(...),
    text: Optional[str] = Form(None),
    file: Optional[UploadFile] = File(None)
):
    all_flashcards = []

    # Extraction and splitting based on the method
    if method == "pdf":
        if not file or not file.filename.endswith(".pdf"):
            raise HTTPException(status_code=400, detail="Please upload a valid PDF file.")
        extracted_text = extract_text_from_pdf(file.file)
    elif method == "pptx":
        if not file or not file.filename.endswith(".pptx"):
            raise HTTPException(status_code=400, detail="Please upload a valid PPTX file.")
        extracted_text = extract_text_from_pptx(file.file)
    elif method == "docx":
        if not file or not file.filename.endswith(".docx"):
            raise HTTPException(status_code=400, detail="Please upload a valid DOCX file.")
        extracted_text = extract_text_from_docx(file.file)
    elif method == "csv":
        if not file or not file.filename.endswith(".csv"):
            raise HTTPException(status_code=400, detail="Please upload a valid CSV file.")
        extracted_text = extract_text_from_csv(file.file)
    elif method == "text":
        if not text:
            raise HTTPException(status_code=400, detail="Please provide valid text input.")
        extracted_text = text
    else:
        raise HTTPException(status_code=400, detail="Invalid method specified.")

    # Split the document into manageable chunks
    chunks = RecursiveCharacterTextSplitter(chunk_size=750, chunk_overlap=100).split_text(extracted_text)

    # Select the appropriate prompt based on the flashcard type
    if type == "type-I":
        structured_llm = llm.with_structured_output(FlashcardSet)
        prompt = normal_prompt
    elif type == "type-II":
        structured_llm = llm.with_structured_output(ClozeDeletionFlashcardSet)
        prompt = cloze_prompt
    else:
        raise HTTPException(status_code=400, detail="Invalid type specified.")

    # Generate flashcards for the extracted chunks
    for chunk in chunks:
        if chunk.strip():
            try:
                response = structured_llm.invoke(prompt.format(chunk=chunk))
                if 'flashcards' in response:
                    all_flashcards.extend(response['flashcards'])
            except Exception as e:
                print(f"Error generating flashcards for chunk")

    return {"flashcards": all_flashcards}

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000, debug=True)
